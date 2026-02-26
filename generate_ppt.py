from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Wissen Seat Booking System"
    subtitle.text = "A modern, intelligent desk allocation solution\nFinal Submission - 4 PM"

    # Slide 2: Problem Statement & Constraints
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "System Requirements & Constraints"
    tf = body_shape.text_frame
    
    tf.text = "Capacity Constraint: 50 Total Office Seats vs 80 Employees"
    
    p = tf.add_paragraph()
    p.text = "Workforce Distribution:"
    p = tf.add_paragraph()
    p.text = "- 10 Squads (8 members each) = 80 Total Members"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Split into Batches: Batch 1 (40) vs Batch 2 (40)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Designated Attendance Days:"
    p = tf.add_paragraph()
    p.text = "- Batch 1: Mon, Tue, Wed"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Batch 2: Thu, Fri"
    p.level = 1

    # Slide 3: Allocation Logic & Business Rules
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Smart Allocation Logic"
    tf = body_shape.text_frame
    
    tf.text = "Designated Seats:"
    p = tf.add_paragraph()
    p.text = "- Auto-assigned to respective batches on their designated days."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Can be booked at any time on the given day."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Buffer / Floating Seats (10 Seats):"
    p = tf.add_paragraph()
    p.text = "- Essential for cross-batch collaboration (e.g., Batch 1 employee working on Thursday)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Time Gate: Can ONLY be booked after 3 PM the day prior to ensure fair forecasting."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Dynamic Resourcing:"
    p = tf.add_paragraph()
    p.text = "- When an employee releases a 'Designated' seat they aren't using, it is automatically converted into a 'Floating' seat for that day, expanding the pool for others instantly."
    p.level = 1

    # Slide 4: Key Strengths & Tech Stack (USP)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Technology Stack & USP (Unique Selling Proposition)"
    tf = body_shape.text_frame
    
    tf.text = "Modern Tech Stack:"
    p = tf.add_paragraph()
    p.text = "- Frontend: React (Vite), Tailwind CSS, Shadcn UI Components (Sonner Toasts)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Backend: Node.js (Express), Server-less configuration"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Database: Firebase Firestore Add-Ons"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "USPs (Unique Selling Propositions):"
    p = tf.add_paragraph()
    p.text = "1. Socket.io Real-Time Synchronization: Instantly updates 'Available Buffer Seats' globally for all connected employees the millisecond someone confirms or cancels a booking. Say goodbye to race conditions."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Security & Dynamics: Git scraped of hardcoded API keys. Dynamic profile page allowing you to simulate different batch/squad roles instantly."
    p.level = 1

    prs.save('Wissen_Seat_Booking_Presentation.pptx')
    print("Presentation saved successfully as Wissen_Seat_Booking_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
